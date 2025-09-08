import os
import sys
import tempfile
import logging
from pathlib import Path
from typing import Optional, Tuple, Dict, Any
from datetime import datetime
import base64
import platform

# 跨平台文档处理库
try:
    # PPT处理
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors

    # Word处理
    from docx import Document

    # 图片处理
    from PIL import Image as PILImage

    # 跨平台Word转换
    import mammoth
    from weasyprint import HTML, CSS

    # Windows特定的Word转换（可选）
    if platform.system().lower() == 'windows':
        try:
            from docx2pdf import convert as docx2pdf_convert
            DOCX2PDF_AVAILABLE = True
        except ImportError:
            DOCX2PDF_AVAILABLE = False
    else:
        DOCX2PDF_AVAILABLE = False

    # 不再使用COM接口，完全使用纯Python实现
    COM_AVAILABLE = False

    print("跨平台文档转换库加载成功")

except ImportError as e:
    print(f"导入库失败: {e}")
    print("请安装必需的依赖:")
    print("pip install python-pptx reportlab python-docx mammoth weasyprint pillow")
    if platform.system().lower() == 'windows':
        print("pip install docx2pdf  # Windows可选")

logger = logging.getLogger(__name__)


class DocumentConverterService:
    """
    跨平台文档转换服务
    支持Word和PPT文档转换为PDF，不依赖Microsoft Office
    """

    def __init__(self):
        self.temp_dir = tempfile.gettempdir()
        self.supported_formats = {
            'ppt': self.convert_ppt_to_pdf_cross_platform,
            'pptx': self.convert_ppt_to_pdf_cross_platform,
            'doc': self.convert_word_to_pdf_cross_platform,
            'docx': self.convert_word_to_pdf_cross_platform
        }

        # 转换进度回调
        self.progress_callback = None

        # 确保输出目录存在
        os.makedirs(self.temp_dir, exist_ok=True)
        logger.info(f"文档转换器初始化完成，输出目录: {self.temp_dir}")
        logger.info(f"当前平台: {platform.system()}")
        logger.info(f"docx2pdf可用: {DOCX2PDF_AVAILABLE}")

    def set_progress_callback(self, callback):
        """设置进度回调函数"""
        self.progress_callback = callback

    def _update_progress(self, stage: str, progress: int, message: str):
        """更新转换进度"""
        if self.progress_callback:
            self.progress_callback({
                'stage': stage,
                'progress': progress,
                'message': message,
                'timestamp': datetime.now().isoformat()
            })

    def convert_ppt_to_pdf_cross_platform(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """
        跨平台PPT转PDF转换

        Args:
            input_path: PPT文件路径
            output_path: 输出PDF路径

        Returns:
            转换结果信息
        """
        try:
            logger.info(f"开始PPT转PDF: {input_path} -> {output_path}")

            # 进度: 开始解析PPT
            self._update_progress('parsing', 10, '正在解析PPT文件...')

            # 加载PPT文件
            presentation = Presentation(input_path)
            slide_count = len(presentation.slides)

            logger.info(f"PPT包含 {slide_count} 张幻灯片")

            # 进度: PPT解析完成
            self._update_progress('parsed', 20, f'PPT解析完成，共{slide_count}张幻灯片')

            # 创建PDF文档
            doc = SimpleDocTemplate(output_path, pagesize=A4)
            story = []
            styles = getSampleStyleSheet()

            # 创建标题样式
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=12,
                textColor=colors.darkblue
            )

            # 创建内容样式
            content_style = ParagraphStyle(
                'CustomContent',
                parent=styles['Normal'],
                fontSize=12,
                spaceAfter=6,
                leftIndent=20
            )

            # 进度: 开始处理幻灯片
            self._update_progress('processing', 30, '正在处理幻灯片内容...')

            # 处理每张幻灯片
            for i, slide in enumerate(presentation.slides, 1):
                logger.info(f"处理第 {i} 张幻灯片")

                # 更新进度
                slide_progress = 30 + int((i / slide_count) * 50)  # 30-80%
                self._update_progress('processing', slide_progress, f'正在处理第{i}/{slide_count}张幻灯片...')

                # 添加幻灯片标题
                story.append(Paragraph(f"幻灯片 {i}", title_style))
                story.append(Spacer(1, 12))

                # 提取文本内容
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                # 添加文本内容
                if slide_text:
                    for text in slide_text:
                        if text:
                            story.append(Paragraph(text, content_style))
                            story.append(Spacer(1, 6))
                else:
                    story.append(Paragraph("(此幻灯片无文本内容)", content_style))

                # 添加分页符（除了最后一页）
                if i < slide_count:
                    story.append(Spacer(1, 20))
                    story.append(Paragraph("─" * 50, styles['Normal']))
                    story.append(Spacer(1, 20))

            # 进度: 生成PDF
            self._update_progress('generating', 85, '正在生成PDF文件...')

            # 生成PDF
            doc.build(story)

            # 进度: 完成
            self._update_progress('complete', 100, 'PPT转PDF完成！')

            logger.info(f"PPT转PDF完成: {slide_count} 张幻灯片")

            return {
                'method': 'python-pptx + reportlab',
                'pages': slide_count,
                'slides_processed': slide_count
            }

        except Exception as e:
            logger.error(f"PPT转PDF失败: {str(e)}")
            raise

    def convert_word_to_pdf_cross_platform(self, input_path: str, output_path: str) -> Dict[str, Any]:
        """
        跨平台Word转PDF转换

        Args:
            input_path: Word文件路径
            output_path: 输出PDF路径

        Returns:
            转换结果信息
        """
        try:
            logger.info(f"开始Word转PDF: {input_path} -> {output_path}")

            # 进度: 开始转换
            self._update_progress('starting', 10, '正在准备Word转换...')

            # 方法1: 尝试使用docx2pdf (Windows)
            if DOCX2PDF_AVAILABLE and platform.system().lower() == 'windows':
                try:
                    self._update_progress('converting', 30, '使用docx2pdf转换中...')
                    docx2pdf_convert(input_path, output_path)
                    if os.path.exists(output_path):
                        self._update_progress('complete', 100, 'Word转PDF完成！')
                        logger.info("使用docx2pdf转换成功")
                        return {
                            'method': 'docx2pdf',
                            'pages': self._count_pdf_pages(output_path)
                        }
                except Exception as e:
                    logger.warning(f"docx2pdf转换失败，尝试备用方法: {e}")

            # 方法2: 使用mammoth + weasyprint (跨平台)
            self._update_progress('parsing', 20, '正在解析Word文档...')
            logger.info("使用mammoth + weasyprint进行转换")

            # 读取Word文档
            self._update_progress('parsing', 40, '正在读取Word文档内容...')
            with open(input_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html_content = result.value

                if result.messages:
                    logger.warning(f"转换警告: {result.messages}")

            # 进度: HTML转换完成
            self._update_progress('html_ready', 60, '文档内容解析完成，正在生成HTML...')

            # 创建完整的HTML文档
            full_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <meta charset="utf-8">
                <style>
                    body {{
                        font-family: Arial, sans-serif;
                        line-height: 1.6;
                        margin: 40px;
                        color: #333;
                    }}
                    h1, h2, h3, h4, h5, h6 {{
                        color: #2c3e50;
                        margin-top: 20px;
                        margin-bottom: 10px;
                    }}
                    p {{
                        margin-bottom: 10px;
                    }}
                    table {{
                        border-collapse: collapse;
                        width: 100%;
                        margin: 20px 0;
                    }}
                    th, td {{
                        border: 1px solid #ddd;
                        padding: 8px;
                        text-align: left;
                    }}
                    th {{
                        background-color: #f2f2f2;
                    }}
                </style>
            </head>
            <body>
                {html_content}
            </body>
            </html>
            """

            # 进度: 生成PDF
            self._update_progress('generating', 80, '正在生成PDF文件...')

            # 转换为PDF
            HTML(string=full_html).write_pdf(output_path)

            # 进度: 完成
            self._update_progress('complete', 100, 'Word转PDF完成！')

            logger.info("Word转PDF完成")

            return {
                'method': 'mammoth + weasyprint',
                'pages': self._count_pdf_pages(output_path)
            }

        except Exception as e:
            logger.error(f"Word转PDF失败: {str(e)}")
            raise

    def _count_pdf_pages(self, pdf_path: str) -> int:
        """
        统计PDF页数

        Args:
            pdf_path: PDF文件路径

        Returns:
            页数
        """
        try:
            # 这里可以使用PyPDF2或其他库来统计页数
            # 简单起见，返回1
            return 1
        except:
            return 1

    def convert_file(self, input_path: str, output_format: str = 'pdf') -> Dict[str, Any]:
        """
        统一的文件转换接口

        Args:
            input_path: 输入文件路径
            output_format: 输出格式，目前只支持pdf

        Returns:
            转换结果字典
        """
        try:
            # 验证输入文件
            if not os.path.exists(input_path):
                raise FileNotFoundError(f"输入文件不存在: {input_path}")

            # 获取文件信息
            file_path = Path(input_path)
            file_extension = file_path.suffix.lower().lstrip('.')
            file_name = file_path.stem

            logger.info(f"开始转换文件: {input_path}")
            logger.info(f"文件类型: {file_extension}")

            # 检查是否支持该格式
            if file_extension not in self.supported_formats:
                raise ValueError(f"不支持的文件格式: {file_extension}")

            # 生成输出文件路径
            output_filename = f"{file_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            output_path = os.path.join(self.temp_dir, output_filename)

            # 执行转换
            converter_func = self.supported_formats[file_extension]
            conversion_result = converter_func(input_path, output_path)

            # 验证输出文件
            if not os.path.exists(output_path):
                raise RuntimeError("转换失败：输出文件未生成")

            # 获取文件大小
            output_size = os.path.getsize(output_path)

            result = {
                'success': True,
                'input_file': input_path,
                'output_file': output_path,
                'output_size': output_size,
                'file_type': file_extension,
                'conversion_method': conversion_result.get('method', 'unknown'),
                'pages': conversion_result.get('pages', 0),
                'timestamp': datetime.now().isoformat(),
                'message': '转换成功'
            }

            logger.info(f"转换完成: {output_path} ({output_size} bytes)")
            return result

        except Exception as e:
            logger.error(f"转换失败: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'input_file': input_path,
                'timestamp': datetime.now().isoformat(),
                'message': f'转换失败: {str(e)}'
            }

    def convert_base64_file(self, base64_data: str, file_extension: str) -> Dict[str, Any]:
        """
        转换base64编码的文件

        Args:
            base64_data: base64编码的文件数据
            file_extension: 文件扩展名

        Returns:
            转换结果
        """
        try:
            # 解码base64数据
            file_data = base64.b64decode(base64_data)

            # 创建临时输入文件
            temp_input = tempfile.NamedTemporaryFile(
                suffix=f'.{file_extension}',
                delete=False
            )

            with open(temp_input.name, 'wb') as f:
                f.write(file_data)

            # 转换文件
            result = self.convert_file(temp_input.name)

            # 如果转换成功，读取输出文件并转换为base64
            if result['success']:
                with open(result['output_file'], 'rb') as f:
                    pdf_data = f.read()
                    result['pdf_base64'] = base64.b64encode(pdf_data).decode('utf-8')

            # 清理临时文件
            try:
                os.unlink(temp_input.name)
                if result['success']:
                    os.unlink(result['output_file'])
            except:
                pass

            return result

        except Exception as e:
            logger.error(f"base64文件转换失败: {str(e)}")
            return {
                'success': False,
                'error': str(e),
                'message': f'base64文件转换失败: {str(e)}'
            }

    def convert_word_to_pdf(self, word_file_path: str, output_path: Optional[str] = None) -> Tuple[bool, str, str]:
        """
        将Word文档转换为PDF（使用纯Python实现）

        Args:
            word_file_path: Word文档路径
            output_path: 输出PDF路径，如果为None则自动生成

        Returns:
            Tuple[bool, str, str]: (成功标志, PDF路径, 错误信息)
        """
        try:
            if not os.path.exists(word_file_path):
                return False, "", f"Word文件不存在: {word_file_path}"

            # 生成输出路径
            if output_path is None:
                base_name = Path(word_file_path).stem
                output_path = os.path.join(self.temp_dir, f"{base_name}_converted.pdf")

            logger.info(f"开始转换Word文档（纯Python）: {word_file_path} -> {output_path}")

            # 使用跨平台方法转换
            result = self.convert_word_to_pdf_cross_platform(word_file_path, output_path)

            return True, output_path, "转换成功"

        except Exception as e:
            logger.error(f"Word转换服务异常: {str(e)}")
            return False, "", f"转换服务异常: {str(e)}"
    
    def convert_ppt_to_pdf(self, ppt_file_path: str, output_path: Optional[str] = None) -> Tuple[bool, str, str]:
        """
        将PPT文档转换为PDF（使用纯Python实现）

        Args:
            ppt_file_path: PPT文档路径
            output_path: 输出PDF路径，如果为None则自动生成

        Returns:
            Tuple[bool, str, str]: (成功标志, PDF路径, 错误信息)
        """
        try:
            if not os.path.exists(ppt_file_path):
                return False, "", f"PPT文件不存在: {ppt_file_path}"

            # 生成输出路径
            if output_path is None:
                base_name = Path(ppt_file_path).stem
                output_path = os.path.join(self.temp_dir, f"{base_name}_converted.pdf")

            logger.info(f"开始转换PPT文档（纯Python）: {ppt_file_path} -> {output_path}")

            # 使用跨平台方法转换
            result = self.convert_ppt_to_pdf_cross_platform(ppt_file_path, output_path)

            return True, output_path, "转换成功"

        except Exception as e:
            logger.error(f"PPT转换服务异常: {str(e)}")
            return False, "", f"转换服务异常: {str(e)}"
    
    def convert_document(self, file_path: str, output_path: Optional[str] = None) -> Tuple[bool, str, str]:
        """
        通用文档转换方法
        根据文件扩展名自动选择转换方法
        
        Args:
            file_path: 文档路径
            output_path: 输出PDF路径
            
        Returns:
            Tuple[bool, str, str]: (成功标志, PDF路径, 错误信息)
        """
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext in ['.doc', '.docx']:
                return self.convert_word_to_pdf(file_path, output_path)
            elif file_ext in ['.ppt', '.pptx']:
                return self.convert_ppt_to_pdf(file_path, output_path)
            else:
                return False, "", f"不支持的文件格式: {file_ext}"
                
        except Exception as e:
            logger.error(f"文档转换异常: {str(e)}")
            return False, "", f"文档转换异常: {str(e)}"
    
    def cleanup_temp_files(self, max_age_hours: int = 24):
        """
        清理临时文件
        
        Args:
            max_age_hours: 文件最大保留时间（小时）
        """
        try:
            import time
            current_time = time.time()
            temp_path = Path(self.temp_dir)
            
            for file_path in temp_path.glob("*_converted.pdf"):
                file_age = current_time - file_path.stat().st_mtime
                if file_age > max_age_hours * 3600:  # 转换为秒
                    try:
                        file_path.unlink()
                        logger.info(f"清理临时文件: {file_path}")
                    except Exception as e:
                        logger.warning(f"清理文件失败 {file_path}: {str(e)}")
                        
        except Exception as e:
            logger.error(f"清理临时文件异常: {str(e)}")


# 创建全局实例
document_converter = DocumentConverterService()
